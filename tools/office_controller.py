"""
tools/office_controller.py — Automatización de Microsoft Office para Alisha IA.

Funciones:
  - read_word(path)              → lee texto de un .docx
  - create_word(path, text)      → crea documento Word con contenido
  - read_excel(path, sheet)      → lee celdas de un .xlsx
  - read_powerpoint(path)        → lee texto de cada diapositiva
  - save_and_close_word(path)    → guarda y cierra documento Word

Principio fail-silent: toda excepción retorna string descriptivo, sin raise.
"""
from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Optional

logger = logging.getLogger("OfficeController")
if not logger.handlers:
    logging.basicConfig(level=logging.INFO, format="[%(name)s] %(levelname)s: %(message)s")

# ── Dependencias opcionales ────────────────────────────────────────────────────
try:
    from docx import Document as _DocxDocument
    _DOCX_OK = True
except ImportError:
    _DocxDocument = None  # type: ignore
    _DOCX_OK = False
    logger.warning("python-docx no instalado. Instalar con: pip install python-docx")

try:
    import openpyxl as _openpyxl
    _OPENPYXL_OK = True
except ImportError:
    _openpyxl = None  # type: ignore
    _OPENPYXL_OK = False
    logger.warning("openpyxl no instalado. Instalar con: pip install openpyxl")

try:
    from pptx import Presentation as _Presentation
    _PPTX_OK = True
except ImportError:
    _Presentation = None  # type: ignore
    _PPTX_OK = False
    logger.warning("python-pptx no instalado. Instalar con: pip install python-pptx")


# ══════════════════════════════════════════════════════════════════════════════
# WORD (Req 3.6, 3.7)
# ══════════════════════════════════════════════════════════════════════════════

def read_word(path: str) -> str:
    """
    Lee el contenido de texto de un archivo Word (.docx).

    Args:
        path: Ruta al archivo .docx

    Returns:
        Texto completo del documento, o mensaje de error.
    """
    if not _DOCX_OK:
        return "Error: python-docx no instalado. Instalar con: pip install python-docx"

    try:
        ruta = Path(path)
        if not ruta.exists():
            return f"Error: El archivo no existe: {path}"
        if not ruta.suffix.lower() in (".docx", ".doc"):
            return f"Error: Formato no soportado: {ruta.suffix}. Solo .docx"

        doc = _DocxDocument(str(ruta))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        logger.info("Word leído: %s (%d párrafos)", ruta.name, len(paragraphs))
        return text if text else "(Documento vacío)"

    except Exception as e:
        logger.warning("Error al leer Word '%s': %s", path, e)
        return f"Error al leer el documento Word: {e}"


def create_word(path: str, text: str, title: str = "") -> str:
    """
    Crea un nuevo documento Word con el texto proporcionado.

    Args:
        path:  Ruta donde guardar el .docx
        text:  Contenido del documento
        title: Título opcional (se agrega como encabezado)

    Returns:
        Mensaje de éxito o error.
    """
    if not _DOCX_OK:
        return "Error: python-docx no instalado. Instalar con: pip install python-docx"

    try:
        ruta = Path(path)
        ruta.parent.mkdir(parents=True, exist_ok=True)

        doc = _DocxDocument()

        # Título opcional
        if title:
            doc.add_heading(title, level=1)

        # Agregar texto párrafo por párrafo
        for linea in text.split("\n"):
            doc.add_paragraph(linea)

        doc.save(str(ruta))
        logger.info("Word creado: %s", ruta.name)
        return f"✓ Documento Word creado: {ruta}"

    except Exception as e:
        logger.warning("Error al crear Word '%s': %s", path, e)
        return f"Error al crear el documento Word: {e}"


def save_and_close_word(path: str) -> str:
    """
    Guarda y cierra un documento Word abierto (via win32com si está disponible).

    Args:
        path: Ruta del documento a cerrar

    Returns:
        Mensaje de resultado.
    """
    try:
        import win32com.client as win32
        word = win32.GetActiveObject("Word.Application")
        for doc in word.Documents:
            if Path(doc.FullName).resolve() == Path(path).resolve():
                doc.Save()
                doc.Close()
                return f"✓ Documento guardado y cerrado: {Path(path).name}"
        return f"No se encontró el documento abierto: {Path(path).name}"
    except ImportError:
        # Fallback: solo guardar con python-docx (no puede cerrar la app)
        return f"win32com no disponible. El documento fue guardado pero no se puede cerrar automáticamente."
    except Exception as e:
        return f"Error al guardar/cerrar Word: {e}"


# ══════════════════════════════════════════════════════════════════════════════
# EXCEL (Req 3.8)
# ══════════════════════════════════════════════════════════════════════════════

def read_excel(path: str, sheet_name: Optional[str] = None) -> list[list]:
    """
    Lee el contenido de un archivo Excel (.xlsx).

    Args:
        path:       Ruta al archivo .xlsx
        sheet_name: Nombre de la hoja (None = primera hoja)

    Returns:
        Lista de filas, cada fila es una lista de valores.
        Retorna lista con mensaje de error si falla.
    """
    if not _OPENPYXL_OK:
        return [["Error: openpyxl no instalado. Instalar con: pip install openpyxl"]]

    try:
        ruta = Path(path)
        if not ruta.exists():
            return [[f"Error: El archivo no existe: {path}"]]
        if not ruta.suffix.lower() in (".xlsx", ".xls", ".xlsm"):
            return [[f"Error: Formato no soportado: {ruta.suffix}"]]

        wb = _openpyxl.load_workbook(str(ruta), read_only=True, data_only=True)

        # Seleccionar hoja
        if sheet_name and sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = []
        for row in ws.iter_rows(values_only=True):
            # Filtrar filas completamente vacías
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell for cell in row_data):
                rows.append(row_data)

        wb.close()
        logger.info("Excel leído: %s (%d filas)", ruta.name, len(rows))
        return rows if rows else [["(Hoja vacía)"]]

    except Exception as e:
        logger.warning("Error al leer Excel '%s': %s", path, e)
        return [[f"Error al leer el archivo Excel: {e}"]]


# ══════════════════════════════════════════════════════════════════════════════
# POWERPOINT (Req 3.9)
# ══════════════════════════════════════════════════════════════════════════════

def read_powerpoint(path: str) -> list[str]:
    """
    Lee el texto de cada diapositiva de un archivo PowerPoint (.pptx).

    Args:
        path: Ruta al archivo .pptx

    Returns:
        Lista de strings, uno por diapositiva con su contenido de texto.
    """
    if not _PPTX_OK:
        return ["Error: python-pptx no instalado. Instalar con: pip install python-pptx"]

    try:
        ruta = Path(path)
        if not ruta.exists():
            return [f"Error: El archivo no existe: {path}"]
        if not ruta.suffix.lower() in (".pptx", ".ppt"):
            return [f"Error: Formato no soportado: {ruta.suffix}"]

        prs = _Presentation(str(ruta))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_content = f"[Diapositiva {i}]\n" + "\n".join(texts)
            slides_text.append(slide_content)

        logger.info("PowerPoint leído: %s (%d diapositivas)", ruta.name, len(slides_text))
        return slides_text if slides_text else ["(Presentación vacía)"]

    except Exception as e:
        logger.warning("Error al leer PowerPoint '%s': %s", path, e)
        return [f"Error al leer la presentación: {e}"]


# ══════════════════════════════════════════════════════════════════════════════
# ABRIR APLICACIONES OFFICE (Req 3.10)
# ══════════════════════════════════════════════════════════════════════════════

def open_office_app(app: str, file_path: str = "") -> str:
    """
    Abre una aplicación de Office (Word, Excel, PowerPoint).

    Args:
        app:       "word", "excel", "powerpoint"
        file_path: Ruta opcional al archivo a abrir

    Returns:
        Mensaje de resultado.
    """
    import subprocess

    app_map = {
        "word":        "WINWORD.EXE",
        "excel":       "EXCEL.EXE",
        "powerpoint":  "POWERPNT.EXE",
    }

    exe = app_map.get(app.lower())
    if not exe:
        return f"Aplicación no reconocida: {app}. Opciones: word, excel, powerpoint"

    # Buscar el ejecutable
    try:
        from config.settings import _resolver_app
        ruta_exe = _resolver_app(app.lower())
    except Exception:
        ruta_exe = exe

    try:
        cmd = [ruta_exe]
        if file_path and Path(file_path).exists():
            cmd.append(file_path)
        subprocess.Popen(cmd, shell=False)
        logger.info("Abriendo %s%s", app, f": {file_path}" if file_path else "")
        return f"✓ Abriendo {app.capitalize()}{f': {Path(file_path).name}' if file_path else ''}."
    except Exception as e:
        return f"Error al abrir {app}: {e}"
